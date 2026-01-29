import streamlit as st
import pandas as pd
import plotly.express as px
import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re

# --- 1. CẤU HÌNH TRANG ---
st.set_page_config(page_title="Hệ Thống Quản Trị Đội Xe", page_icon="🚘", layout="wide")

st.markdown("""
<style>
    .block-container {padding-top: 1rem; padding-bottom: 3rem;}
    .kpi-card {
        background-color: white; border-radius: 12px; padding: 20px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05); border: 1px solid #f0f2f6;
        height: 100%; display: flex; flex-direction: column; justify-content: space-between;
        min-height: 160px;
    }
    .kpi-header { display: flex; justify-content: space-between; align-items: center; margin-bottom: 10px; }
    .kpi-title { font-size: 14px; color: #6c757d; font-weight: 700; text-transform: uppercase; }
    .kpi-icon { font-size: 20px; background: #f8f9fa; padding: 8px; border-radius: 8px; }
    .kpi-value { font-size: 32px; font-weight: 800; color: #212529; margin: 0; }
    .kpi-formula { font-size: 12px; color: #888; font-style: italic; margin-top: 10px; border-top: 1px dashed #eee; padding-top: 5px; }
    .progress-bg { background-color: #e9ecef; border-radius: 4px; height: 6px; width: 100%; margin: 8px 0; overflow: hidden; }
    .progress-fill { height: 100%; border-radius: 4px; transition: width 0.5s ease-in-out; }
</style>
""", unsafe_allow_html=True)

# --- 2. HÀM XỬ LÝ DỮ LIỆU (ĐÃ FIX LỖI TRÙNG CỘT BU) ---
@st.cache_data
def load_data_final(file):
    try:
        # 1. Đọc file
        if file.name.endswith('.csv'):
            df = pd.read_csv(file)
        else:
            df = pd.read_excel(file, engine='openpyxl')

        # 2. Chuẩn hóa tên cột
        df.columns = [str(c).strip().replace('\n', ' ') for c in df.columns]

        # --- FIX LỖI QUAN TRỌNG: XỬ LÝ CỘT 'BU' GỐC ---
        # File gốc có cột 'BU' (chứa mã số) và 'Bộ phận' (chứa tên). 
        # Ta cần xóa hoặc đổi tên cột 'BU' gốc trước khi map 'Bộ phận' vào 'BU' để tránh trùng lặp.
        if 'BU' in df.columns:
            df = df.rename(columns={'BU': 'BU_Code_Original'}) # Đổi tên cột mã số đi để tránh trùng

        # 3. MAPPING CỘT
        rename_map = {
            'Ngày Tháng Năm': 'Ngày khởi hành',
            'Biển số xe': 'Biển số xe',
            'Tên tài xế': 'Tên tài xế',
            'Người sử dụng xe': 'Người sử dụng xe',
            'Lộ trình': 'Lộ trình',
            'Giờ khởi hành': 'Giờ khởi hành',
            'Giờ kết thúc': 'Giờ kết thúc',
            'Công Ty': 'Công ty',
            'Bộ phận': 'BU',             # Map Bộ phận (Tên) vào biến BU để vẽ biểu đồ
            'Site': 'Location',
            'Cost center': 'Cost Center',
            'Tổng chi phí': 'Chi phí'
        }
        df = df.rename(columns=rename_map)

        # 4. LÀM SẠCH DỮ LIỆU CƠ BẢN
        if 'Ngày khởi hành' in df.columns:
            df['Ngày khởi hành'] = pd.to_datetime(df['Ngày khởi hành'], errors='coerce')
            df = df.dropna(subset=['Ngày khởi hành']) 

            # Fix lỗi năm 2026
            mask_error = (df['Ngày khởi hành'].dt.month > 6) & (df['Ngày khởi hành'].dt.year == 2026)
            if mask_error.any():
                df.loc[mask_error, 'Ngày khởi hành'] = df.loc[mask_error, 'Ngày khởi hành'].apply(lambda x: x.replace(month=1))

        # 5. XỬ LÝ SỐ LIỆU CHI PHÍ
        if 'Chi phí' in df.columns:
            df['Chi phí'] = df['Chi phí'].astype(str).str.replace(',', '').str.replace('.', '', regex=False)
            df['Chi phí'] = pd.to_numeric(df['Chi phí'], errors='coerce').fillna(0)
        else:
            df['Chi phí'] = 0

        # 6. TÍNH TOÁN THỜI GIAN
        df['Start'] = pd.to_datetime(df['Ngày khởi hành'].astype(str) + ' ' + df['Giờ khởi hành'].astype(str), errors='coerce')
        df['End'] = pd.to_datetime(df['Ngày khởi hành'].astype(str) + ' ' + df['Giờ kết thúc'].astype(str), errors='coerce')
        
        mask_overnight = df['End'] < df['Start']
        df.loc[mask_overnight, 'End'] += pd.Timedelta(days=1)
        
        df['Duration'] = (df['End'] - df['Start']).dt.total_seconds() / 3600
        df['Tháng'] = df['Start'].dt.strftime('%Y-%m')
        
        # 7. PHÂN LOẠI XE & CÁC CỘT KHÁC
        def normalize_plate(plate):
            if not isinstance(plate, str): return ""
            return re.sub(r'[^A-Z0-9]', '', plate.upper())
        
        if 'Biển số xe' in df.columns:
            df['Biển_Clean'] = df['Biển số xe'].apply(normalize_plate)

        if 'Phân Loại Xe' not in df.columns:
             df['Phân Loại Xe'] = df['Tên tài xế'].apply(lambda x: 'Xe Nội bộ' if pd.notna(x) and str(x).strip() != '' else 'Xe Vãng lai')
        
        # Điền Unknown cho các cột thiếu
        for col in ['Cost Center', 'Công ty', 'BU', 'Location']:
            if col not in df.columns: df[col] = 'Unknown'
            df[col] = df[col].fillna('Unknown').astype(str)

        # Fix lỗi KeyError Tình trạng đơn
        if 'Tình trạng đơn yêu cầu' not in df.columns:
            df['Tình trạng đơn yêu cầu'] = 'APPROVED'
        
        # Phân loại phạm vi
        def check_scope(r):
            s = str(r).lower()
            provinces = ['bình dương', 'đồng nai', 'long an', 'bà rịa', 'vũng tàu', 'tây ninh', 'bình phước', 'tiền giang', 'bến tre', 'cần thơ', 'vĩnh long', 'an giang', 'phan thiết', 'mũi né', 'trà vinh', 'bắc ninh', 'hải phòng']
            if any(p in s for p in provinces): return "Đi Tỉnh"
            return "Nội thành"
        
        df['Phạm Vi'] = df['Lộ trình'].apply(check_scope) if 'Lộ trình' in df.columns else 'Unknown'

        return df, {}

    except Exception as e:
        return f"❌ Lỗi xử lý file: {str(e)}", {}

# --- 3. CHART EXPORT ---
def get_chart_img(data, x, y, kind='bar', title='', color='#0078d4'):
    plt.figure(figsize=(7, 4.5))
    if x not in data.columns or y not in data.columns:
        plt.text(0.5, 0.5, 'No Data', ha='center'); img = BytesIO(); plt.savefig(img, format='png'); plt.close(); img.seek(0); return img
    if kind == 'bar': 
        data = data.sort_values(by=x, ascending=True)
        bars = plt.barh(data[y], data[x], color=color); plt.xlabel(x); plt.bar_label(bars, fmt='%g')
    elif kind == 'column': 
        bars = plt.bar(data[y], data[x], color=color); plt.ylabel(x); plt.xticks(rotation=45, ha='right'); plt.bar_label(bars, fmt='%g')
    elif kind == 'pie': 
        plt.pie(data[x], labels=data[y], autopct='%1.1f%%', startangle=90, colors=['#107c10', '#d13438', '#0078d4', '#ffc107', '#8764b8'])
    plt.title(title, pad=15, fontweight='bold', fontsize=12, color='#333'); plt.tight_layout()
    img = BytesIO(); plt.savefig(img, format='png', dpi=120); plt.close(); img.seek(0)
    return img

# --- 4. EXPORT PPTX ---
def export_pptx(kpi, df_comp, df_status, top_users, top_drivers, df_bad_trips, selected_options, chart_prefs, df_scope):
    prs = Presentation()
    
    def add_kpi_shape(slide, left, top, width, height, title, value, sub, color_rgb):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = color_rgb; shape.line.width = Pt(2.5)
        tb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.3))
        tb.text_frame.text = title; tb.text_frame.paragraphs[0].font.bold = True; tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        tb_v = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.4), width - Inches(0.2), Inches(0.5))
        p_v = tb_v.text_frame.paragraphs[0]; p_v.text = str(value); p_v.font.size = Pt(24); p_v.font.bold = True
        tb_s = slide.shapes.add_textbox(left + Inches(0.1), top + height - Inches(0.4), width - Inches(0.2), Inches(0.3))
        p_s = tb_s.text_frame.paragraphs[0]; p_s.text = sub; p_s.font.size = Pt(9); p_s.font.italic = True; p_s.font.color.rgb = RGBColor(150, 150, 150)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BÁO CÁO VẬN HÀNH ĐỘI XE"; slide.placeholders[1].text = f"Cập nhật: {kpi['last_month']}"
    
    slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "TỔNG QUAN HIỆU SUẤT"
    add_kpi_shape(slide, Inches(0.5), Inches(2.5), Inches(1.8), Inches(1.5), "TỔNG CHUYẾN", f"{kpi['trips']}", "Số chuyến", RGBColor(0, 120, 212))
    add_kpi_shape(slide, Inches(2.4), Inches(2.5), Inches(1.8), Inches(1.5), "GIỜ VẬN HÀNH", f"{kpi['hours']:,.0f}", "Tổng giờ", RGBColor(0, 120, 212))
    add_kpi_shape(slide, Inches(4.3), Inches(2.5), Inches(1.8), Inches(1.5), "CÔNG SUẤT", kpi['occupancy_text'], "Mục tiêu >50%", RGBColor(0, 120, 212))
    add_kpi_shape(slide, Inches(6.2), Inches(2.5), Inches(1.8), Inches(1.5), "HOÀN THÀNH", f"{kpi['success_rate']:.1f}%", "Tỷ lệ OK", RGBColor(16, 124, 16))
    add_kpi_shape(slide, Inches(8.1), Inches(2.5), Inches(1.8), Inches(1.5), "HỦY/TỪ CHỐI", f"{kpi['cancel_rate'] + kpi['reject_rate']:.1f}%", "Tỷ lệ Fail", RGBColor(209, 52, 56))

    if "Biểu đồ Tổng quan" in selected_options:
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "PHÂN TÍCH CẤU TRÚC SỬ DỤNG"
        if not df_comp.empty:
            img1 = get_chart_img(df_comp.head(8), 'Value', 'Category', kind=chart_prefs.get('structure', 'bar'), title='Top Đơn Vị')
            slide.shapes.add_picture(img1, Inches(0.5), Inches(1.8), Inches(4.5), Inches(3.5))
        if not df_scope.empty:
            img2 = get_chart_img(df_scope, 'Số lượng', 'Phạm vi', kind=chart_prefs.get('scope', 'pie'), title='Phạm Vi')
            slide.shapes.add_picture(img2, Inches(5.2), Inches(1.8), Inches(4.5), Inches(3.5))

    if "Bảng Xếp Hạng (Top User/Driver)" in selected_options:
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "BẢNG XẾP HẠNG HOẠT ĐỘNG"
        if not top_users.empty:
            img_u = get_chart_img(top_users.head(8), 'Số chuyến', 'Người sử dụng xe', kind=chart_prefs.get('top_user', 'bar'), title='Top Users', color='#8764b8')
            slide.shapes.add_picture(img_u, Inches(0.5), Inches(1.8), Inches(4.5), Inches(3.5))
        if not top_drivers.empty:
            img_d = get_chart_img(top_drivers.head(8), 'Số chuyến', 'Tên tài xế', kind=chart_prefs.get('top_driver', 'bar'), title='Top Drivers', color='#00cc6a')
            slide.shapes.add_picture(img_d, Inches(5.2), Inches(1.8), Inches(4.5), Inches(3.5))

    if "Danh sách Hủy/Từ chối" in selected_options:
        slide = prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text = "CHI TIẾT ĐƠN HỦY / TỪ CHỐI"
        cols_safe = ['Start_Str', 'User', 'Status', 'Note', 'Lý do']
        avail_cols = [c for c in cols_safe if c in df_bad_trips.columns]
        rows, cols = min(len(df_bad_trips)+1, 10), len(avail_cols)
        if cols > 0:
            table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8)).table
            for i, h in enumerate(avail_cols):
                cell = table.cell(0, i); cell.text = h
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0, 120, 212)
            for i, row in enumerate(df_bad_trips.head(9).itertuples(), start=1):
                for j, col_name in enumerate(avail_cols):
                    val = getattr(row, col_name, "")
                    table.cell(i, j).text = str(val)[:30]

    out = BytesIO(); prs.save(out); out.seek(0); return out

# --- 5. GIAO DIỆN CHÍNH ---
st.title("📊 Phước Minh - Hệ Thống Quản Trị & Tối Ưu Hóa Đội Xe")
uploaded_file = st.file_uploader("Upload Excel/CSV", type=['xlsx', 'csv'], label_visibility="collapsed")

if uploaded_file:
    df, report_info = load_data_final(uploaded_file)
    if isinstance(df, str): st.error(df); st.stop()
    
    # SIDEBAR
    with st.sidebar:
        st.header("🗂️ Bộ Lọc Dữ Liệu")
        
        # 1. Thời gian
        min_date, max_date = df['Start'].min().date(), df['Start'].max().date()
        date_range = st.date_input("Thời gian:", (min_date, max_date), min_value=min_date, max_value=max_date)
        
        # 2. Loại Xe
        unique_types = df['Phân Loại Xe'].unique().tolist()
        type_filter = st.multiselect("Loại Xe:", unique_types, default=unique_types)
        
        # LỌC CẤP 1
        df_filtered = df.copy()
        if len(date_range) == 2:
            df_filtered = df_filtered[(df_filtered['Start'].dt.date >= date_range[0]) & (df_filtered['Start'].dt.date <= date_range[1])]
        if type_filter:
            df_filtered = df_filtered[df_filtered['Phân Loại Xe'].isin(type_filter)]

        st.markdown("---")
        st.caption("Drill-down (Theo cấu trúc Công ty):")
        
        # 3. CÔNG TY (Lên đầu tiên)
        comps = ["Tất cả"] + sorted(df_filtered['Công ty'].dropna().unique().tolist())
        sel_comp = st.selectbox("1. Công ty:", comps)
        
        # Lọc theo Công ty trước để Cost Center hiển thị đúng tương ứng
        if sel_comp != "Tất cả": 
            df_filtered = df_filtered[df_filtered['Công ty'] == sel_comp]

        # 4. COST CENTER (Phụ thuộc vào Công ty đã chọn)
        # Chỉ hiện các Cost Center thuộc Công ty (hoặc tất cả nếu chưa chọn Cty)
        available_ccs = sorted(df_filtered['Cost Center'].unique().tolist())
        ccs = ["Tất cả"] + available_ccs
        sel_cc = st.selectbox("2. Cost Center:", ccs)
        
        if sel_cc != "Tất cả": 
            df_filtered = df_filtered[df_filtered['Cost Center'] == sel_cc]

        # 5. Khu vực
        locs = ["Tất cả"] + sorted(df_filtered['Location'].dropna().unique().tolist())
        sel_loc = st.selectbox("3. Khu vực (Site):", locs)
        if sel_loc != "Tất cả": df_filtered = df_filtered[df_filtered['Location'] == sel_loc]
        
        st.write(f"🔍 Đang xem: **{len(df_filtered)}** chuyến")

    if df_filtered.empty: st.warning("Không có dữ liệu."); st.stop()

    # --- KPI CALCULATION ---
    total_trips = len(df_filtered)
    total_hours = df_filtered['Duration'].sum()
    total_cost = df_filtered['Chi phí'].sum()

    # Công suất
    internal_df = df_filtered[df_filtered['Phân Loại Xe'] == 'Xe Nội bộ']
    active_cars_list = internal_df['Biển_Clean'].unique()
    num_active_cars = len(active_cars_list)
    
    if len(date_range) == 2: num_days = (date_range[1] - date_range[0]).days + 1
    else: num_days = 1
        
    capacity_hours = num_active_cars * num_days * 8
    actual_run_hours = internal_df['Duration'].sum()
    occupancy_pct = (actual_run_hours / capacity_hours * 100) if capacity_hours > 0 else 0

    # Tỷ lệ thành công (Safe Mode)
    if 'Tình trạng đơn yêu cầu' in df_filtered.columns:
        counts = df_filtered['Tình trạng đơn yêu cầu'].fillna('Unknown').value_counts()
        suc_rate = ((counts.get('CLOSED', 0) + counts.get('APPROVED', 0)) / total_trips * 100) if total_trips > 0 else 100
        fail_rate = ((counts.get('CANCELED', 0) + counts.get('CANCELLED', 0) + counts.get('REJECTED_BY_ADMIN', 0)) / total_trips * 100) if total_trips > 0 else 0
    else:
        suc_rate = 100.0; fail_rate = 0.0

    # --- UI HIỂN THỊ ---
    st.markdown("### 📈 Tổng Quan Vận Hành")
    debug_mode = st.checkbox("🛠️ Hiển thị công thức tính & Dữ liệu gốc (Debug Mode)")

    cols = st.columns(4)
    with cols[0]:
        st.metric("Tổng Chuyến", f"{total_trips}", delta="Chuyến xe")
        if debug_mode: st.info(f"Đếm số dòng dữ liệu sau lọc: {len(df_filtered)}")
    with cols[1]:
        st.metric("Tổng Chi Phí", f"{total_cost:,.0f} đ", delta="VND")
        if debug_mode: st.info(f"Tổng cột 'Chi phí' (Đã xử lý dấu phẩy).")
    with cols[2]:
        st.metric("Công Suất (Nội bộ)", f"{occupancy_pct:.1f}%")
        if debug_mode: st.warning(f"Công thức: {actual_run_hours:.1f}h chạy / ({num_active_cars} xe * {num_days} ngày * 8h)")
    with cols[3]:
        st.metric("Tổng Giờ Chạy", f"{total_hours:,.0f}h")

    if debug_mode:
        st.markdown("---")
        st.error("🚨 **DỮ LIỆU GỐC:**")
        st.dataframe(df_filtered[['Ngày khởi hành', 'Biển số xe', 'Công ty', 'Cost Center', 'Giờ khởi hành', 'Giờ kết thúc', 'Chi phí']].head(100), use_container_width=True)

    # --- TABS ---
    t1, t2, t3, t4 = st.tabs(["📊 Phân Tích", "🏆 Bảng Xếp Hạng", "📉 Chất Lượng", "⚙️ Đối Soát & Kiểm Tra"])
    chart_prefs = {} 
    kind_map = {"Thanh ngang (Bar)": "bar", "Thanh dọc (Column)": "column", "Tròn (Pie)": "pie"}

    with t1:
        c1, c2 = st.columns([2, 1])
        with c1:
            st.write("#### Phân bổ Cấu trúc (Phòng ban/Công ty)")
            chart_type_struct = st.selectbox("Kiểu biểu đồ:", list(kind_map.keys()), index=0, key="c_struct")
            chart_prefs['structure'] = kind_map[chart_type_struct]
            
            # Logic vẽ biểu đồ thông minh:
            # Nếu Đang chọn "Tất cả" công ty -> Vẽ biểu đồ so sánh các Công ty
            # Nếu Đã chọn 1 công ty cụ thể -> Vẽ biểu đồ so sánh các Phòng ban (BU) bên trong
            if sel_comp == "Tất cả": 
                df_g = df_filtered['Công ty'].value_counts().reset_index()
                df_g.columns = ['Category', 'Value']
                title_c = "Tỷ trọng theo Công Ty"
            else: 
                df_g = df_filtered['BU'].value_counts().reset_index()
                df_g.columns = ['Category', 'Value']
                title_c = f"Tỷ trọng Phòng Ban ({sel_comp})"
            
            if chart_prefs['structure'] == "bar": fig = px.bar(df_g, x='Value', y='Category', orientation='h', text='Value', title=title_c)
            elif chart_prefs['structure'] == "column": fig = px.bar(df_g, x='Category', y='Value', text='Value', title=title_c)
            else: fig = px.pie(df_g, values='Value', names='Category', title=title_c)
            st.plotly_chart(fig, use_container_width=True)
        with c2:
            st.write("#### Phạm vi di chuyển")
            chart_type_scope = st.selectbox("Kiểu biểu đồ Phạm vi:", list(kind_map.keys()), index=2, key="c_scope")
            chart_prefs['scope'] = kind_map[chart_type_scope]
            if 'Phạm Vi' in df_filtered.columns:
                df_sc = df_filtered['Phạm Vi'].value_counts().reset_index(); df_sc.columns = ['Phạm vi', 'Số lượng']
                if chart_prefs['scope'] == "pie": fig_s = px.pie(df_sc, values='Số lượng', names='Phạm vi', hole=0.5)
                elif chart_prefs['scope'] == "bar": fig_s = px.bar(df_sc, x='Số lượng', y='Phạm vi', orientation='h', text='Số lượng')
                else: fig_s = px.bar(df_sc, x='Phạm vi', y='Số lượng', text='Số lượng')
                st.plotly_chart(fig_s, use_container_width=True)

    with t2:
        c_u, c_d = st.columns(2)
        with c_u:
            type_u = st.selectbox("Biểu đồ Top User:", list(kind_map.keys()), index=0, key="c_user")
            chart_prefs['top_user'] = kind_map[type_u]
            top_u = df_filtered.groupby(['Người sử dụng xe', 'Công ty']).size().reset_index(name='Số chuyến').sort_values('Số chuyến', ascending=False).head(10)
            st.write("##### 🥇 Top User")
            if chart_prefs['top_user'] == "bar": fig_u = px.bar(top_u, x='Số chuyến', y='Người sử dụng xe', orientation='h', text='Số chuyến', hover_data=['Công ty'])
            elif chart_prefs['top_user'] == "column": fig_u = px.bar(top_u, x='Người sử dụng xe', y='Số chuyến', text='Số chuyến')
            else: fig_u = px.pie(top_u, values='Số chuyến', names='Người sử dụng xe')
            st.plotly_chart(fig_u, use_container_width=True)
        with c_d:
            type_d = st.selectbox("Biểu đồ Top Driver:", list(kind_map.keys()), index=0, key="c_driver")
            chart_prefs['top_driver'] = kind_map[type_d]
            top_d = df_filtered.groupby(['Tên tài xế', 'Phân Loại Xe']).size().reset_index(name='Số chuyến').sort_values('Số chuyến', ascending=False).head(10)
            st.write("##### 🚘 Top Driver")
            if chart_prefs['top_driver'] == "bar": fig_d = px.bar(top_d, x='Số chuyến', y='Tên tài xế', orientation='h', text='Số chuyến', hover_data=['Phân Loại Xe'])
            elif chart_prefs['top_driver'] == "column": fig_d = px.bar(top_d, x='Tên tài xế', y='Số chuyến', text='Số chuyến')
            else: fig_d = px.pie(top_d, values='Số chuyến', names='Tên tài xế')
            st.plotly_chart(fig_d, use_container_width=True)

    with t3:
        st.write("#### Chi tiết Hủy / Từ chối")
        bad = pd.DataFrame()
        if 'Tình trạng đơn yêu cầu' in df_filtered.columns:
            bad = df_filtered[df_filtered['Tình trạng đơn yêu cầu'].isin(['CANCELED', 'CANCELLED', 'REJECTED_BY_ADMIN'])]
        if not bad.empty: st.dataframe(bad, use_container_width=True)
        else: st.success("Không có chuyến nào bị hủy trong dữ liệu log này.")

    with t4:
        st.subheader("⚙️ Đối Soát Dữ Liệu")
        with st.expander(f"🚗 Danh sách Xe Hoạt Động"):
            unique_cars = df_filtered[['Biển số xe', 'Tên tài xế', 'Phân Loại Xe']].drop_duplicates().sort_values('Biển số xe')
            st.dataframe(unique_cars, use_container_width=True)

    # --- PPTX ---
    st.divider()
    st.subheader("📥 Xuất Báo Cáo PowerPoint")
    c_opt, c_btn = st.columns([2, 1])
    with c_opt:
        pptx_options = st.multiselect("Chọn nội dung Slide:", ["Biểu đồ Tổng quan", "Bảng Xếp Hạng (Top User/Driver)", "Danh sách Hủy/Từ chối"], default=["Biểu đồ Tổng quan", "Bảng Xếp Hạng (Top User/Driver)"])
    with c_btn:
        st.write(""); st.write("")
        last_month_str = "N/A"
        try:
            if not df.empty and 'Tháng' in df.columns:
                valid_months = df['Tháng'].dropna()
                if not valid_months.empty: last_month_str = valid_months.max()
        except: pass

        kpi_data = {'trips': total_trips, 'hours': total_hours, 'occupancy': occupancy_pct, 'occupancy_text': f"{occupancy_pct:.1f}%", 'success_rate': suc_rate, 'cancel_rate': fail_rate, 'reject_rate': 0, 'last_month': last_month_str}
        df_status_exp = pd.DataFrame() # No status chart needed for log
        if sel_comp == "Tất cả": df_comp_exp = df_filtered['Công ty'].value_counts().reset_index(); df_comp_exp.columns=['Category', 'Value']
        else: df_comp_exp = df_filtered['BU'].value_counts().reset_index(); df_comp_exp.columns=['Category', 'Value']
        if 'Phạm Vi' in df_filtered.columns: df_scope_exp = df_filtered['Phạm Vi'].value_counts().reset_index(); df_scope_exp.columns = ['Phạm vi', 'Số lượng']
        else: df_scope_exp = pd.DataFrame(columns=['Phạm vi', 'Số lượng'])
        df_bad_exp = bad.copy() if not bad.empty else pd.DataFrame()
        if not df_bad_exp.empty:
            df_bad_exp['Start_Str'] = df_bad_exp['Start'].dt.strftime('%d/%m')
            df_bad_exp = df_bad_exp.rename(columns={'Người sử dụng xe': 'User', 'Tình trạng đơn yêu cầu': 'Status'})

        pptx_file = export_pptx(kpi_data, df_comp_exp, df_status_exp, top_u, top_d, df_bad_exp, pptx_options, chart_prefs, df_scope_exp)
        st.download_button(label="Tải file .PPTX ngay", data=pptx_file, file_name="Bao_Cao_Van_Hanh_Full.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", type="primary")

else:
    st.info("👋 Vui lòng upload file Excel/CSV dữ liệu mới.")